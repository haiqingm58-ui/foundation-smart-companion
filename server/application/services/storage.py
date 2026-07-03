from __future__ import annotations

import io
from pathlib import Path
from uuid import uuid4

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from ..errors import APIError


ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".png", ".jpg", ".jpeg", ".webp", ".md", ".markdown", ".txt"}
MAX_FILE_SIZE = 25 * 1024 * 1024


def validate_upload(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise APIError(400, "不支持该文件类型", "FILE_TYPE_NOT_ALLOWED")
    if not content:
        raise APIError(400, "文件内容为空", "FILE_EMPTY")
    if len(content) > MAX_FILE_SIZE:
        raise APIError(413, "文件不能超过 25 MB", "FILE_TOO_LARGE")
    return suffix


def save_upload(upload_dir: Path, filename: str, content: bytes) -> Path:
    suffix = validate_upload(filename, content)
    upload_dir.mkdir(parents=True, exist_ok=True)
    destination = (upload_dir / f"{uuid4().hex}{suffix}").resolve()
    if not destination.is_relative_to(upload_dir.resolve()):
        raise APIError(400, "文件路径不安全", "UNSAFE_FILE_PATH")
    destination.write_bytes(content)
    return destination


def extract_text(content: bytes, suffix: str) -> str:
    suffix = suffix.lower()
    if suffix in {".md", ".markdown", ".txt"}:
        return content.decode("utf-8", errors="ignore")
    if suffix == ".pdf":
        return "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(content)).pages)
    if suffix == ".docx":
        return "\n".join(paragraph.text for paragraph in Document(io.BytesIO(content)).paragraphs)
    if suffix == ".pptx":
        presentation = Presentation(io.BytesIO(content))
        return "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
    if suffix == ".xlsx":
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        return "\n".join("\t".join(str(value or "") for value in row) for sheet in workbook.worksheets for row in sheet.iter_rows(values_only=True))
    return ""


def chunk_text(text: str, size: int = 700, overlap: int = 100) -> list[str]:
    compact = "\n".join(line.strip() for line in text.replace("\r", "").splitlines() if line.strip())
    if not compact:
        return []
    chunks = []
    start = 0
    while start < len(compact):
        end = min(len(compact), start + size)
        chunks.append(compact[start:end])
        if end >= len(compact):
            break
        start = max(start + 1, end - overlap)
    return chunks
